"""
Presentation Export Service
Handles exporting presentations to different formats (PDF, PPTX, HTML, etc.)
"""

import os
import json
import base64
import zipfile
from io import BytesIO
from typing import Dict, List, Optional
from django.template.loader import render_to_string
from django.conf import settings


class PresentationExportService:
    """Service for exporting presentations to various formats"""
    
    def __init__(self):
        self.supported_formats = ['pdf', 'pptx', 'html', 'images', 'json']
    
    def export_presentation(self, presentation, export_format: str, 
                          include_notes: bool = True, 
                          high_quality: bool = True) -> Dict:
        """
        Export a presentation to the specified format
        """
        try:
            if export_format not in self.supported_formats:
                raise ValueError(f"Unsupported format: {export_format}")
            
            slides = presentation.slides.all().order_by('slide_number')
            
            if export_format == 'json':
                return self._export_json(presentation, slides, include_notes)
            elif export_format == 'html':
                return self._export_html(presentation, slides, include_notes, high_quality)
            elif export_format == 'pdf':
                return self._export_pdf(presentation, slides, include_notes, high_quality)
            elif export_format == 'pptx':
                return self._export_pptx(presentation, slides, include_notes, high_quality)
            elif export_format == 'images':
                return self._export_images(presentation, slides, high_quality)
            else:
                raise ValueError(f"Export format {export_format} not implemented")
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'file_data': None
            }
    
    def _export_json(self, presentation, slides, include_notes: bool) -> Dict:
        """Export presentation as JSON"""
        data = {
            'presentation': {
                'id': presentation.id,
                'title': presentation.title,
                'description': presentation.description,
                'topic': presentation.topic,
                'target_audience': presentation.target_audience,
                'presentation_type': presentation.presentation_type,
                'theme': presentation.theme,
                'color_scheme': presentation.color_scheme,
                'tone': presentation.tone,
                'created_at': presentation.created_at.isoformat(),
                'slides': []
            }
        }
        
        for slide in slides:
            slide_data = {
                'slide_number': slide.slide_number,
                'title': slide.title,
                'subtitle': slide.subtitle,
                'content': slide.content,
                'slide_type': slide.slide_type,
                'layout': slide.layout,
                'background_color': slide.background_color,
                'text_color': slide.text_color,
                'accent_color': slide.accent_color
            }
            
            if include_notes:
                slide_data['notes'] = slide.notes
            
            # Include slide elements
            elements = []
            for element in slide.elements.all():
                elements.append({
                    'element_type': element.element_type,
                    'position_x': element.position_x,
                    'position_y': element.position_y,
                    'width': element.width,
                    'height': element.height,
                    'content': element.content,
                    'content_data': element.content_data,
                    'styling': {
                        'font_size': element.font_size,
                        'font_weight': element.font_weight,
                        'text_align': element.text_align,
                        'color': element.color,
                        'background': element.background,
                        'border': element.border
                    }
                })
            slide_data['elements'] = elements
            
            data['presentation']['slides'].append(slide_data)
        
        json_str = json.dumps(data, indent=2, ensure_ascii=False)
        encoded_data = base64.b64encode(json_str.encode('utf-8')).decode('utf-8')
        
        return {
            'success': True,
            'file_data': encoded_data,
            'file_size': len(json_str),
            'mime_type': 'application/json',
            'filename': f"{presentation.title}.json"
        }
    
    def _export_html(self, presentation, slides, include_notes: bool, high_quality: bool) -> Dict:
        """Export presentation as HTML"""
        
        # Prepare slide data
        slide_data = []
        for slide in slides:
            slide_info = {
                'slide': slide,
                'elements': slide.elements.all()
            }
            slide_data.append(slide_info)
        
        # Render HTML template
        html_content = render_to_string('presentation_export/html_export.html', {
            'presentation': presentation,
            'slides': slide_data,
            'include_notes': include_notes,
            'high_quality': high_quality
        })
        
        encoded_data = base64.b64encode(html_content.encode('utf-8')).decode('utf-8')
        
        return {
            'success': True,
            'file_data': encoded_data,
            'file_size': len(html_content),
            'mime_type': 'text/html',
            'filename': f"{presentation.title}.html"
        }
    
    def _export_pdf(self, presentation, slides, include_notes: bool, high_quality: bool) -> Dict:
        """Export presentation as PDF"""
        try:
            # First generate HTML
            html_result = self._export_html(presentation, slides, include_notes, high_quality)
            if not html_result['success']:
                return html_result
            
            # Decode HTML content
            html_content = base64.b64decode(html_result['file_data']).decode('utf-8')
            
            # Try to use weasyprint for PDF generation (if available)
            try:
                import weasyprint
                pdf_bytes = weasyprint.HTML(string=html_content).write_pdf()
                encoded_data = base64.b64encode(pdf_bytes).decode('utf-8')
                
                return {
                    'success': True,
                    'file_data': encoded_data,
                    'file_size': len(pdf_bytes),
                    'mime_type': 'application/pdf',
                    'filename': f"{presentation.title}.pdf"
                }
            except ImportError:
                # Fallback to basic PDF-style HTML
                pdf_html = self._generate_pdf_html(presentation, slides, include_notes)
                encoded_data = base64.b64encode(pdf_html.encode('utf-8')).decode('utf-8')
                
                return {
                    'success': True,
                    'file_data': encoded_data,
                    'file_size': len(pdf_html),
                    'mime_type': 'text/html',
                    'filename': f"{presentation.title}_pdf.html"
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'file_data': None
            }
    
    def _export_pptx(self, presentation, slides, include_notes: bool, high_quality: bool) -> Dict:
        """Export presentation as PowerPoint (PPTX)"""
        try:
            # Try to use python-pptx library (if available)
            try:
                from pptx import Presentation as PPTXPresentation
                from pptx.util import Inches
                
                prs = PPTXPresentation()
                
                # Set presentation properties
                prs.core_properties.title = presentation.title
                prs.core_properties.author = presentation.user.get_full_name() or presentation.user.username
                prs.core_properties.subject = presentation.description
                
                for slide in slides:
                    # Add slide layout
                    slide_layout = prs.slide_layouts[0] if slide.slide_type == 'title' else prs.slide_layouts[1]
                    pptx_slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    if hasattr(pptx_slide.shapes, 'title') and slide.title:
                        pptx_slide.shapes.title.text = slide.title
                    
                    # Add content
                    if slide.content and len(pptx_slide.placeholders) > 1:
                        content_placeholder = pptx_slide.placeholders[1]
                        content_placeholder.text = slide.content
                    
                    # Add speaker notes
                    if include_notes and slide.notes:
                        notes_slide = pptx_slide.notes_slide
                        notes_slide.notes_text_frame.text = slide.notes
                
                # Save to bytes
                pptx_bytes = BytesIO()
                prs.save(pptx_bytes)
                pptx_bytes.seek(0)
                
                encoded_data = base64.b64encode(pptx_bytes.getvalue()).decode('utf-8')
                
                return {
                    'success': True,
                    'file_data': encoded_data,
                    'file_size': len(pptx_bytes.getvalue()),
                    'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'filename': f"{presentation.title}.pptx"
                }
                
            except ImportError:
                # Fallback to HTML-based export
                return self._export_html(presentation, slides, include_notes, high_quality)
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'file_data': None
            }
    
    def _export_images(self, presentation, slides, high_quality: bool) -> Dict:
        """Export presentation as image files in a ZIP"""
        try:
            # Create ZIP file in memory
            zip_buffer = BytesIO()
            
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                # Add a simple text file with presentation info
                info_content = f"""Presentation: {presentation.title}
Description: {presentation.description or 'No description'}
Slides: {len(slides)}
Theme: {presentation.theme}
Created: {presentation.created_at}

Slides:
"""
                for slide in slides:
                    info_content += f"\nSlide {slide.slide_number}: {slide.title or 'Untitled'}\n"
                    if slide.content:
                        info_content += f"Content: {slide.content[:100]}...\n"
                
                zip_file.writestr('presentation_info.txt', info_content)
                
                # For now, add HTML versions of each slide as individual files
                for slide in slides:
                    slide_html = self._generate_slide_html(slide, high_quality)
                    filename = f"slide_{slide.slide_number:02d}_{slide.title or 'untitled'}.html"
                    # Clean filename
                    filename = "".join(c for c in filename if c.isalnum() or c in "._-").replace(" ", "_")
                    zip_file.writestr(filename, slide_html)
            
            zip_buffer.seek(0)
            encoded_data = base64.b64encode(zip_buffer.getvalue()).decode('utf-8')
            
            return {
                'success': True,
                'file_data': encoded_data,
                'file_size': len(zip_buffer.getvalue()),
                'mime_type': 'application/zip',
                'filename': f"{presentation.title}_slides.zip"
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'file_data': None
            }
    
    def _generate_pdf_html(self, presentation, slides, include_notes: bool) -> str:
        """Generate PDF-optimized HTML"""
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>{presentation.title}</title>
            <style>
                @page {{
                    size: A4 landscape;
                    margin: 1in;
                }}
                body {{
                    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                    margin: 0;
                    padding: 0;
                    line-height: 1.6;
                }}
                .slide {{
                    page-break-after: always;
                    padding: 20px;
                    min-height: 80vh;
                    display: flex;
                    flex-direction: column;
                }}
                .slide:last-child {{
                    page-break-after: auto;
                }}
                .slide-header {{
                    border-bottom: 3px solid #3B82F6;
                    padding-bottom: 10px;
                    margin-bottom: 20px;
                }}
                .slide-title {{
                    font-size: 28px;
                    font-weight: bold;
                    color: #1F2937;
                    margin: 0;
                }}
                .slide-subtitle {{
                    font-size: 18px;
                    color: #6B7280;
                    margin: 5px 0 0 0;
                }}
                .slide-content {{
                    flex-grow: 1;
                    font-size: 16px;
                    color: #374151;
                    white-space: pre-wrap;
                }}
                .slide-notes {{
                    margin-top: 20px;
                    padding-top: 15px;
                    border-top: 1px solid #E5E7EB;
                    font-style: italic;
                    color: #6B7280;
                    font-size: 14px;
                }}
                .presentation-title {{
                    text-align: center;
                    font-size: 36px;
                    color: #1F2937;
                    margin-bottom: 40px;
                }}
            </style>
        </head>
        <body>
            <div class="slide">
                <div class="presentation-title">{presentation.title}</div>
                <div style="text-align: center; font-size: 18px; color: #6B7280;">
                    {presentation.description or ''}
                </div>
                <div style="text-align: center; margin-top: 40px;">
                    <strong>Generated by IntelliHub AI</strong><br>
                    {presentation.created_at.strftime('%B %d, %Y')}
                </div>
            </div>
        """
        
        for slide in slides:
            html_content += f"""
            <div class="slide">
                <div class="slide-header">
                    <h1 class="slide-title">{slide.title or f'Slide {slide.slide_number}'}</h1>
                    {f'<div class="slide-subtitle">{slide.subtitle}</div>' if slide.subtitle else ''}
                </div>
                <div class="slide-content">{slide.content or ''}</div>
                {f'<div class="slide-notes"><strong>Speaker Notes:</strong><br>{slide.notes}</div>' if include_notes and slide.notes else ''}
            </div>
            """
        
        html_content += """
        </body>
        </html>
        """
        
        return html_content
    
    def _generate_slide_html(self, slide, high_quality: bool) -> str:
        """Generate HTML for a single slide"""
        return f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>Slide {slide.slide_number}</title>
            <style>
                body {{
                    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                    margin: 0;
                    padding: 40px;
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    min-height: 100vh;
                    display: flex;
                    align-items: center;
                    justify-content: center;
                }}
                .slide-container {{
                    background: white;
                    padding: 60px;
                    border-radius: 20px;
                    box-shadow: 0 20px 40px rgba(0,0,0,0.1);
                    max-width: 1000px;
                    width: 100%;
                    aspect-ratio: 16/9;
                }}
                .slide-title {{
                    font-size: 48px;
                    font-weight: bold;
                    color: #1F2937;
                    margin-bottom: 20px;
                    text-align: center;
                }}
                .slide-content {{
                    font-size: 24px;
                    color: #374151;
                    line-height: 1.8;
                    white-space: pre-wrap;
                }}
            </style>
        </head>
        <body>
            <div class="slide-container">
                <h1 class="slide-title">{slide.title or f'Slide {slide.slide_number}'}</h1>
                <div class="slide-content">{slide.content or ''}</div>
            </div>
        </body>
        </html>
        """


# Service instance
export_service = PresentationExportService()


def export_presentation_to_format(presentation, export_format: str, 
                                 include_notes: bool = True, 
                                 high_quality: bool = True) -> Dict:
    """
    Main function to export a presentation to a specific format
    """
    return export_service.export_presentation(
        presentation, export_format, include_notes, high_quality
    )


def get_supported_export_formats() -> List[Dict]:
    """Get list of supported export formats with descriptions"""
    return [
        {
            'id': 'pdf',
            'name': 'PDF Document',
            'description': 'Professional PDF suitable for printing and sharing',
            'icon': 'fas fa-file-pdf',
            'color': 'text-red-400'
        },
        {
            'id': 'html',
            'name': 'HTML Presentation',
            'description': 'Interactive web presentation that can be opened in any browser',
            'icon': 'fas fa-code',
            'color': 'text-blue-400'
        },
        {
            'id': 'pptx',
            'name': 'PowerPoint (.pptx)',
            'description': 'Native PowerPoint format for editing in Microsoft Office',
            'icon': 'fas fa-file-powerpoint',
            'color': 'text-orange-400'
        },
        {
            'id': 'json',
            'name': 'JSON Data',
            'description': 'Raw presentation data for developers and integrations',
            'icon': 'fas fa-file-code',
            'color': 'text-green-400'
        },
        {
            'id': 'images',
            'name': 'Image Files (ZIP)',
            'description': 'Individual slide images in a compressed archive',
            'icon': 'fas fa-images',
            'color': 'text-purple-400'
        }
    ]